"""
Text extraction for uploaded files (PDF, DOCX, PPTX, CSV, XLSX, HTML, images, Markdown, plain text).
"""

from __future__ import annotations

import base64
import logging
import os

logger = logging.getLogger(__name__)

ALLOWED_MIME_TYPES = {
    "application/pdf": "PDF",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "Word",
    "application/msword": "Word",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "PowerPoint",
    "application/vnd.ms-powerpoint": "PowerPoint",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "Excel",
    "application/vnd.ms-excel": "Excel",
    "text/csv": "CSV",
    "text/plain": "Text",
    "text/markdown": "Markdown",
    "text/x-markdown": "Markdown",
    "text/html": "HTML",
    "image/jpeg": "Image",
    "image/jpg": "Image",
    "image/png": "Image",
    "image/gif": "Image",
    "image/webp": "Image",
}


async def extract_text(file_path: str, mime_type: str) -> str:
    """Extract plain text from a file based on its mime type."""

    if mime_type == "application/pdf":
        return _extract_pdf(file_path)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return _extract_docx(file_path)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return _extract_pptx(file_path)
    elif mime_type == "text/csv":
        return _extract_csv(file_path)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return _extract_xlsx(file_path)
    elif mime_type == "text/html":
        return _extract_html(file_path)
    elif mime_type.startswith("image/"):
        return await _extract_image(file_path, mime_type)
    elif mime_type in ("text/markdown", "text/x-markdown"):
        return _read_text(file_path)
    elif mime_type.startswith("text/"):
        return _read_text(file_path)
    else:
        logger.warning("Unsupported mime type for extraction: %s", mime_type)
        return _read_text(file_path)


def _read_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_pdf(file_path: str) -> str:
    try:
        import fitz  # pymupdf
    except ImportError:
        logger.error("pymupdf not installed — cannot extract PDF text")
        return ""

    text_parts = []
    with fitz.open(file_path) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts)


def _extract_docx(file_path: str) -> str:
    try:
        import docx
    except ImportError:
        logger.error("python-docx not installed — cannot extract DOCX text")
        return ""

    doc = docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed — cannot extract PPTX text")
        return ""

    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"--- Slide {i} ---"]
        if slide.shapes.title and slide.shapes.title.text:
            parts.append(slide.shapes.title.text)
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text != (slide.shapes.title.text if slide.shapes.title else ""):
                    parts.append(text)
        slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text)


def _extract_csv(file_path: str) -> str:
    try:
        import pandas as pd
    except ImportError:
        logger.error("pandas not installed — cannot extract CSV text")
        return _read_text(file_path)

    df = pd.read_csv(file_path, nrows=200)
    return df.to_markdown(index=False)


def _extract_xlsx(file_path: str) -> str:
    try:
        import openpyxl
    except ImportError:
        logger.error("openpyxl not installed — cannot extract XLSX text")
        return ""

    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        ws = wb.active
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= 200:
                break
            rows.append([str(c) if c is not None else "" for c in row])
        wb.close()

        if not rows:
            return ""

        # Convert to markdown table
        header = rows[0]
        lines = ["| " + " | ".join(header) + " |"]
        lines.append("| " + " | ".join("---" for _ in header) + " |")
        for row in rows[1:]:
            # Pad row to match header length
            padded = row + [""] * (len(header) - len(row))
            lines.append("| " + " | ".join(padded[:len(header)]) + " |")
        return "\n".join(lines)
    except Exception:
        # Fall back to pandas
        try:
            import pandas as pd
            df = pd.read_excel(file_path, nrows=200)
            return df.to_markdown(index=False)
        except Exception as exc:
            logger.error("Failed to extract XLSX: %s", exc)
            return ""


def _extract_html(file_path: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        logger.error("beautifulsoup4 not installed — cannot extract HTML text")
        return _read_text(file_path)

    raw = _read_text(file_path)
    soup = BeautifulSoup(raw, "html.parser")
    text = soup.get_text(separator="\n")
    # Strip blank lines
    return "\n".join(line for line in text.splitlines() if line.strip())


async def _extract_image(file_path: str, mime_type: str) -> str:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        logger.warning("OPENAI_API_KEY not set — cannot OCR image")
        return ""

    base_url = os.getenv("OPENAI_BASE_URL", "https://api.openai.com")

    try:
        import httpx

        with open(file_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode("utf-8")

        # Map mime for data URI
        data_uri = f"data:{mime_type};base64,{image_data}"

        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                f"{base_url}/v1/chat/completions",
                headers={"Authorization": f"Bearer {api_key}"},
                json={
                    "model": "gpt-4o-mini",
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": "Extract all visible text from this image. Return only the text, no commentary.",
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {"url": data_uri},
                                },
                            ],
                        }
                    ],
                    "max_tokens": 4096,
                },
            )
            resp.raise_for_status()
            return resp.json()["choices"][0]["message"]["content"]
    except Exception as exc:
        logger.error("Image OCR failed: %s", exc)
        return ""
